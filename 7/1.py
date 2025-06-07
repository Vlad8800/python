from pptx import Presentation
from pptx.util import Inches, Pt # type: ignore
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])
title_shape = slide_1.shapes.title
left = Inches(1)
top = Inches(1)
slide_1.shapes.add_picture("/mnt/data/logo.png", left, top, height=Inches(2))
title_shape = slide_1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.5))
tf = title_shape.text_frame
tf.text = "Презентація логотипу персонального бренду"

# Slide 2: Ідея та концепція логотипу
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
slide_2.shapes.title.text = "Ідея та концепція логотипу"
content = slide_2.placeholders[1]
content.text = (
    "Назва бренду: NextWeb Studio\n"
    "Сфера: Розробка веб-сайтів та електронної комерції\n"
    "Ідея: Сучасність, надійність та динамізм у веб-розробці\n"
    "Унікальність: Логотип включає стилізовані ініціали 'WD' та елемент веб-браузера"
)

# Slide 3: Кольори та шрифти
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_3.shapes.title.text = "Кольори та шрифти"
content = slide_3.placeholders[1]
content.text = (
    "Кольорова гама:\n"
    "- Синій: довіра, технологічність\n"
    "- Білий: чистота, простота\n"
    "- Чорний: професіоналізм\n\n"
    "Шрифт:\n"
    "- Без засічок, строгий — символізує сучасність і надійність"
)

# Slide 4: Варіанти використання
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
slide_4.shapes.title.text = "Варіанти використання логотипу"
content = slide_4.placeholders[1]
content.text = (
    "- Візитки\n"
    "- Соціальні мережі\n"
    "- Веб-сайти та підписи до електронної пошти\n"
    "- Презентації та документи\n"
)

# Slide 5: Приклади застосування
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])
slide_5.shapes.title.text = "Приклади застосування логотипу"
slide_5.shapes.add_picture("/mnt/data/logo.png", Inches(1), Inches(1.2), height=Inches(2.5))

# Збереження файлу
pptx_file_path = "/mnt/data/NextWeb_Studio_Logo_Presentation.pptx"
prs.save(pptx_file_path)

pptx_file_path
